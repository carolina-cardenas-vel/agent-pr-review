#!/usr/bin/env python3
"""
Generate PowerPoint presentation from presentation outline.
Requires: pip install python-pptx

Usage:
    python3 generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define color scheme
ACCENT_COLOR = RGBColor(0, 102, 204)  # Blue
DARK_COLOR = RGBColor(51, 51, 51)     # Dark gray
LIGHT_COLOR = RGBColor(242, 242, 242) # Light gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE

def add_content_slide(prs, title, points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    # Add separator line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(9.5), Inches(1.2))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(2)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

def generate_presentation():
    """Generate the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "GitHub Copilot PR Review Agent", 
                   "Automating Code Reviews with AI\nBuilding My First Custom Agent\n3 months with Copilot")
    
    # Slide 2: Problem
    add_content_slide(prs, "Project Overview - The Problem", [
        "Code reviews are time-consuming and repetitive",
        "Teams use both GitHub (repos) and Azure DevOps (work items)",
        "Need to verify requirements are actually implemented",
        "Security, performance, and test coverage often get missed",
        "Developers want actionable feedback, not just criticisms",
        "",
        "The Gap: No tool connects ADO requirements to GitHub code"
    ])
    
    # Slide 3: Solution
    add_content_slide(prs, "Project Overview - The Solution", [
        "A VS Code Agent that automatically reviews PRs",
        "Analyzes code against Azure DevOps acceptance criteria",
        "Runs 6 specialized review passes",
        "Posts inline comments directly to PRs",
        "",
        "In plain English: Type review command → Copilot analyzes → Comments appear",
        "",
        "Outcome: Faster reviews, better code quality, fewer bugs shipped"
    ])
    
    # Slide 4: Architecture
    add_content_slide(prs, "Architecture Overview", [
        "1. User triggers review in VS Code Agent Mode",
        "2. Agent fetches ADO work item (requirements)",
        "3. Agent fetches PR files and diffs (code changes)",
        "4. Agent runs 6 analysis passes",
        "5. Agent posts comments to GitHub PR",
        "",
        "Key: VS Code Agent + GitHub Copilot + Two MCP Servers",
        "No custom backend needed—everything runs locally in VS Code"
    ])
    
    # Slide 5: The 6 Passes
    add_content_slide(prs, "The 6 Review Passes", [
        "Pass 1: ADO Compliance — Is each requirement implemented?",
        "Pass 2: Functional Bugs — Null checks, error handling, race conditions",
        "Pass 3: Security — SQL injection, XSS, hardcoded secrets (OWASP Top 10)",
        "Pass 4: Performance — N+1 queries, inefficient loops, unbounded queries",
        "Pass 5: Test Coverage — New functions and branches have tests",
        "Pass 6: Maintainability — Duplication, magic values, tech debt"
    ])
    
    # Slide 6: Example Comment
    add_content_slide(prs, "Example Comment", [
        "🔴 Critical: SQL Injection vulnerability (Pass 3: Security)",
        "File: src/db/queries.ts | Line: 23",
        "Issue: Customer ID directly interpolated into SQL string",
        "Risk: Attacker can inject SQL to steal customer data",
        "Fix: Use parameterized queries",
        "",
        "Color-coded: 🔴 Critical (must fix) | 🟡 Warning (should fix) | 🔵 Suggestion"
    ])
    
    # Slide 7: AI Leverage - Discovery
    add_content_slide(prs, "How I Leveraged AI - Discovery Phase", [
        "Started with a problem: 'I need a PR review tool'",
        "Asked Copilot: 'Create objective.md for this project'",
        "Reviewed the output together, asked follow-up questions",
        "Asked: 'What's the implementation approach?'",
        "Copilot suggested using MCP servers (before I knew they existed!)",
        "Asked: 'Create a detailed plan'",
        "",
        "Key insight: Copilot helped me think through the problem"
    ])
    
    # Slide 8: AI Leverage - Implementation
    add_content_slide(prs, "How I Leveraged AI - Implementation Phase", [
        "Copilot generated: .vscode/mcp.json, instructions, plan, docs",
        "What I didn't have to write: No backend code, no API integration, no OAuth handling",
        "Copilot generated entire 6-pass review logic",
        "",
        "Total setup time: ~2 hours from idea to working agent",
        "Most of that: Understanding MCP servers",
        "Actual coding: ~30 minutes",
        "",
        "Result: Complete, working agent without infrastructure"
    ])
    
    # Slide 9: AI Leverage - Refinement
    add_content_slide(prs, "How I Leveraged AI - Quality & Refinement", [
        "Asked: 'Check this for bugs' → Copilot reviewed instructions",
        "Asked: 'Simplify this pass' → Made review logic clearer",
        "Asked: 'Add examples' → Generated example comments and workflows",
        "Asked: 'Update docs' → Created README, GETTING_STARTED guide",
        "",
        "The power: No trial-and-error coding, no debugging loops",
        "Copilot understood full context from start"
    ])
    
    # Slide 10: Key Learnings - MCP
    add_content_slide(prs, "Key Learnings - MCP Servers", [
        "MCP (Model Context Protocol) servers are pre-built integrations",
        "Official ones exist for major platforms (GitHub, ADO, etc.)",
        "CHECK WHAT'S AVAILABLE BEFORE BUILDING CUSTOM",
        "They handle OAuth automatically",
        "No API tokens to manage or expose",
        "",
        "Key takeaway: Would have taken weeks to build custom integrations"
    ])
    
    # Slide 11: Key Learnings - Agent Mode
    add_content_slide(prs, "Key Learnings - Agent Mode", [
        "VS Code Agent Mode runs Copilot with specific instructions",
        "Activates on trigger text patterns",
        "Automatically accesses MCP tools",
        "Instructions file is incredibly powerful",
        "No UI coding needed—all in prompts",
        "",
        "Key takeaway: Like giving Copilot a job description"
    ])
    
    # Slide 12: Key Learnings - Collaboration
    add_content_slide(prs, "Key Learnings - AI-Human Collaboration", [
        "✓ Start with problem, not solution",
        "✓ Review each artifact and iterate",
        "✓ Use AI for thinking, not just coding",
        "✓ Build on previous outputs (compound effect)",
        "",
        "Each answer informed the next question",
        "Small adjustments unlock better results",
        "Best results from reviewing Copilot's own work"
    ])
    
    # Slide 13: Key Learnings - Prompting
    add_content_slide(prs, "Key Learnings - Prompting Matters", [
        "Context is everything—explain what problem you're solving",
        "Specific examples help—show desired output format",
        "Edge cases matter—mention constraints",
        "Iteration beats perfection—good → great → excellent",
        "",
        "Best results from:",
        "• Asking Copilot to review its own work",
        "• Providing feedback on what worked/didn't work",
        "• Being specific about use case"
    ])
    
    # Slide 14: Gotchas - Auth
    add_content_slide(prs, "Gotchas - MCP Server Authentication", [
        "What I thought: 'Need to manage GitHub PATs and ADO tokens'",
        "What actually happened: MCP servers use OAuth (browser login)",
        "",
        "No tokens stored locally",
        "Session-based authentication",
        "First-time setup: Just sign in once, that's it",
        "",
        "Lesson: Read MCP docs first! Authentication can be simpler"
    ])
    
    # Slide 15: Gotchas - Line Numbers
    add_content_slide(prs, "Gotchas - Diff Line Numbers Are Confusing", [
        "What I thought: 'Just post comments at any line in the file'",
        "What actually happened:",
        "",
        "GitHub API requires line numbers from the diff itself",
        "Can't comment on lines outside changed section",
        "Had to map logic to nearest changed line",
        "Took trial-and-error to get right",
        "",
        "Lesson: Test with real examples early"
    ])
    
    # Slide 16: Gotchas - Prompt Clarity
    add_content_slide(prs, "Gotchas - Vague Instructions = Confused Agent", [
        "What I thought: 'Just tell Copilot to review the code'",
        "What actually happened:",
        "",
        "Generic instructions produced generic comments",
        "Had to write detailed step-by-step logic",
        "6 passes needed to be explicit in order",
        "Formatting rules had to be exact",
        "",
        "Lesson: Be extremely specific with agent instructions"
    ])
    
    # Slide 17: Gotchas - Token Limits
    add_content_slide(prs, "Gotchas - Context Gets Expensive", [
        "What I thought: 'Just fetch all the context'",
        "What actually happened:",
        "",
        "Large PRs with full file context can exceed token limits",
        "Copilot has limits on how much context to process",
        "Had to be selective about which files to fully fetch",
        "",
        "Lesson: Token limits are real constraints"
    ])
    
    # Slide 18: Do Differently - Architecture
    add_content_slide(prs, "What I'd Do Differently - Architecture", [
        "Start with a real PR, not in a vacuum",
        "Test assumptions against actual code early",
        "",
        "Limit scope initially—maybe start with 3 passes, not 6",
        "Add complexity incrementally",
        "Easier to enhance than simplify",
        "",
        "Focus on one problem first (compliance OR security OR performance)"
    ])
    
    # Slide 19: Do Differently - Process
    add_content_slide(prs, "What I'd Do Differently - Process", [
        "Document assumptions explicitly",
        "('We only review TypeScript PRs' would save debugging)",
        "",
        "Test the feedback loop early",
        "Post real comments to test PR, get user feedback",
        "",
        "Build a feedback mechanism",
        "How do we know if comments are helpful?",
        "Track which findings devs actually fix"
    ])
    
    # Slide 20: Do Differently - Copilot Usage
    add_content_slide(prs, "What I'd Do Differently - Using Copilot", [
        "Ask more verification questions: 'Is this the right approach?'",
        "Request specific trade-offs: 'Fast vs. deep analysis?'",
        "",
        "Prototype with real MCP responses early",
        "Ask: 'What does GitHub MCP actually return?'",
        "Test against real API responses",
        "",
        "Current design based on guesses—test early"
    ])
    
    # Slide 21: Do Differently - Documentation
    add_content_slide(prs, "What I'd Do Differently - Documentation", [
        "Architecture diagrams earlier—visual flow catches issues",
        "Use cases & examples upfront (3 real PRs we want to review)",
        "Design from use cases, not theory",
        "",
        "Failure scenarios: 'What if MCP times out?'",
        "'What if ADO work item is closed?'",
        "Build robustness proactively"
    ])
    
    # Slide 22: Key Wins
    add_content_slide(prs, "Key Wins", [
        "✅ MCP servers just worked (OAuth seamless)",
        "✅ Copilot understood full scope (no context re-explaining)",
        "✅ Rapid iteration possible (idea → working agent: 2 hours)",
        "✅ No infrastructure needed (everything runs in VS Code)",
        "",
        "Zero deployment complexity",
        "No trial-and-error or debugging loops"
    ])
    
    # Slide 23: Impact & Next Steps
    add_content_slide(prs, "Impact & Next Steps", [
        "Current: Review any PR against any ADO work item",
        "6 specialized analysis passes, inline comments with fixes",
        "",
        "Immediate next: Deploy and test with real team",
        "Measure impact (bugs caught, time saved)",
        "",
        "Future: Auto-run on every PR, custom rules, metrics dashboard,",
        "multi-language support, continuous learning"
    ])
    
    # Slide 24: Key Takeaways
    add_content_slide(prs, "Lessons for Your Copilot Journey", [
        "1. Use Copilot for thinking, not just coding",
        "   Best results from collaboration, not just code generation",
        "",
        "2. Understand your tools (MCP, APIs, platforms)",
        "   80% work is understanding integrations, 20% is implementation",
        "",
        "3. Iterate with real feedback loops",
        "   Theory beats practice initially; real data changes everything"
    ])
    
    # Slide 25: Q&A
    add_content_slide(prs, "Questions?", [
        "How does it handle edge cases?",
        "Can it integrate with other work item systems?",
        "What's the accuracy of findings?",
        "How do you customize review criteria?",
        "Can it run on every PR automatically?",
        "",
        "Thank you!"
    ])
    
    # Save
    prs.save('GitHub_Copilot_PR_Review_Agent.pptx')
    print("✅ Presentation created: GitHub_Copilot_PR_Review_Agent.pptx")

if __name__ == '__main__':
    try:
        generate_presentation()
    except ImportError:
        print("❌ Missing 'python-pptx' library")
        print("\nInstall it with:")
        print("   pip install python-pptx")
        print("\nThen run this script again:")
        print("   python3 generate_presentation.py")
